# main.py
import os
import io
import re
import zipfile
import tempfile
from typing import List, Tuple

from fastapi import FastAPI, UploadFile, File, HTTPException, Request
from fastapi.responses import JSONResponse, StreamingResponse

# File parsing
from docx import Document  # python-docx
from reportlab.pdfgen import canvas  # reportlab
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm

# Optional deps you should add to requirements.txt:
# pypdf
# python-pptx
try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

# OpenAI (openai>=1.x)
from openai import OpenAI


app = FastAPI()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

MODEL = os.getenv("MODEL", "gpt-4o-mini")  # change if you want
MAX_FILES = int(os.getenv("MAX_FILES", "10"))
MAX_MB_PER_FILE = int(os.getenv("MAX_MB_PER_FILE", "25"))
MAX_TOTAL_MB = int(os.getenv("MAX_TOTAL_MB", "50"))

# Optional simple protection so only n8n can call you:
# Set BACKEND_SHARED_SECRET in Render, then n8n adds header X-Backend-Secret
SHARED_SECRET = os.getenv("BACKEND_SHARED_SECRET", "").strip()


@app.get("/health")
def health():
    return {"ok": True}


def _require_secret(request: Request):
    if not SHARED_SECRET:
        return
    got = request.headers.get("X-Backend-Secret", "")
    if got != SHARED_SECRET:
        raise HTTPException(status_code=401, detail="Unauthorized")


def _bytes_limit_check(files_bytes: List[Tuple[str, bytes]]):
    total = sum(len(b) for _, b in files_bytes)
    if total > MAX_TOTAL_MB * 1024 * 1024:
        raise HTTPException(status_code=413, detail=f"Total upload too large (> {MAX_TOTAL_MB}MB)")


def _read_uploadfile(u: UploadFile) -> bytes:
    data = u.file.read()
    if len(data) > MAX_MB_PER_FILE * 1024 * 1024:
        raise HTTPException(status_code=413, detail=f"{u.filename} too large (> {MAX_MB_PER_FILE}MB)")
    return data


def _safe_filename(name: str) -> str:
    name = name or "file"
    name = re.sub(r"[^a-zA-Z0-9._-]+", "_", name)
    return name[:120]


def extract_text_from_pdf(data: bytes) -> str:
    if PdfReader is None:
        raise HTTPException(status_code=500, detail="PDF support missing. Add 'pypdf' to requirements.txt")
    reader = PdfReader(io.BytesIO(data))
    parts = []
    for i, page in enumerate(reader.pages):
        txt = page.extract_text() or ""
        txt = txt.strip()
        if txt:
            parts.append(txt)
    return "\n\n".join(parts).strip()


def extract_text_from_docx(data: bytes) -> str:
    doc = Document(io.BytesIO(data))
    parts = []
    for p in doc.paragraphs:
        t = (p.text or "").strip()
        if t:
            parts.append(t)
    return "\n".join(parts).strip()


def extract_text_from_pptx(data: bytes) -> str:
    if Presentation is None:
        raise HTTPException(status_code=500, detail="PPTX support missing. Add 'python-pptx' to requirements.txt")
    prs = Presentation(io.BytesIO(data))
    parts = []
    for si, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    slide_text.append(t)
        # Speaker notes (optional)
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_text.append(f"(Notes) {notes}")
        except Exception:
            pass

        if slide_text:
            parts.append(f"Slide {si}:\n" + "\n".join(slide_text))
    return "\n\n".join(parts).strip()


def extract_text(filename: str, data: bytes) -> str:
    fn = (filename or "").lower()
    if fn.endswith(".pdf"):
        return extract_text_from_pdf(data)
    if fn.endswith(".docx"):
        return extract_text_from_docx(data)
    if fn.endswith(".pptx"):
        return extract_text_from_pptx(data)
    if fn.endswith(".txt") or fn.endswith(".md"):
        try:
            return data.decode("utf-8", errors="ignore").strip()
        except Exception:
            return ""
    raise HTTPException(status_code=400, detail=f"Unsupported file type: {filename}")


def call_ai_make_notes(source_text: str) -> str:
    if not os.getenv("OPENAI_API_KEY"):
        raise HTTPException(status_code=500, detail="OPENAI_API_KEY not set")

    # Keep it simple: one prompt that returns markdown notes.
    system = (
    "You are an expert pharmacy educator creating exam-ready study material. "
    "You write clearly, concisely, and accurately for pharmacy students. "
    "You strictly use only the provided source material and do not invent facts."
)

    )
    user = (
    "Create TWO SEPARATE SECTIONS from the source material below.\n\n"

    "================================\n"
    "SECTION 1: STUDY NOTES (PHARMACY)\n"
    "================================\n\n"

    "Purpose:\n"
    "- For a student studying this topic for the FIRST TIME.\n\n"

    "Rules:\n"
    "- Do NOT remove or omit anything required to meet the learning objectives.\n"
    "- Organize content primarily by learning objective.\n"
    "- Be concise but explanatory.\n"
    "- Consolidate repeated ideas.\n"
    "- Do NOT invent facts or add external knowledge.\n"
    "- If something is not covered in the source, say: \"Not covered in these slides\".\n\n"

    "Required structure:\n"
    "1. Title\n"
    "2. Learning Objectives (as given)\n"
    "3. Big Picture Overview (5–6 bullets explaining how the topic fits together)\n"
    "4. Core Notes\n"
    "   - One section per learning objective\n"
    "   - Clear headings and bullet points\n"
    "5. Diagrams & Flowcharts\n"
    "   - Use simple text-based diagrams where helpful (mechanisms, pathways, comparisons)\n"
    "6. Additional Information\n"
    "   - Include content not essential for learning objectives\n"
    "7. Key Terms (simple, pharmacy-relevant definitions)\n\n"

    "========================================\n"
    "SECTION 2: RAPID REVIEW (PHARMACY EXAM)\n"
    "========================================\n\n"

    "Purpose:\n"
    "- For LAST-MINUTE EXAM REVISION.\n\n"

    "Rules:\n"
    "- Extremely concise.\n"
    "- Bullet points only.\n"
    "- No explanations unless essential.\n"
    "- Focus on drug names, mechanisms, indications, contraindications, key interactions,\n"
    "  and NHS/NICE service facts IF present in the source.\n"
    "- Do NOT include anything not in the source.\n\n"

    "Required structure:\n"
    "- Key Facts\n"
    "- High-Yield Drug Points\n"
    "- Common Exam Traps / Confusions\n"
    "- 5–10 Exam-Style Questions (short answer or MCQ style)\n\n"

    "================================\n"
    "SOURCE MATERIAL\n"
    "================================\n"
    f"{source_text}"
)


    resp = client.chat.completions.create(
        model=MODEL,
        messages=[
            {"role": "system", "content": system},
            {"role": "user", "content": user},
        ],
        temperature=0.2,
    )
    out = (resp.choices[0].message.content or "").strip()
    if not out:
        raise HTTPException(status_code=500, detail="AI returned empty notes")
    return out


def markdown_to_docx(md: str) -> bytes:
    doc = Document()

    for line in md.splitlines():
        line = line.rstrip()
        if not line.strip():
            continue

        # Headings
        if line.startswith("### "):
            doc.add_heading(line[4:].strip(), level=3)
            continue
        if line.startswith("## "):
            doc.add_heading(line[3:].strip(), level=2)
            continue
        if line.startswith("# "):
            doc.add_heading(line[2:].strip(), level=1)
            continue

        # Bullets
        if line.lstrip().startswith(("-", "*")):
            text = line.lstrip()[1:].strip()
            p = doc.add_paragraph(text, style="List Bullet")
            continue

        # Numbered
        if re.match(r"^\d+\.\s+", line.strip()):
            text = re.sub(r"^\d+\.\s+", "", line.strip())
            doc.add_paragraph(text, style="List Number")
            continue

        # Plain
        doc.add_paragraph(line.strip())

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def markdown_to_pdf(md: str) -> bytes:
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=A4)
    width, height = A4

    x = 2 * cm
    y = height - 2 * cm
    line_height = 14

    def draw_line(text: str, is_heading=False):
        nonlocal y
        if y < 2 * cm:
            c.showPage()
            y = height - 2 * cm
        if is_heading:
            c.setFont("Helvetica-Bold", 13)
        else:
            c.setFont("Helvetica", 11)
        # Simple wrap
        max_chars = 95
        chunks = [text[i:i + max_chars] for i in range(0, len(text), max_chars)] or [""]
        for ch in chunks:
            if y < 2 * cm:
                c.showPage()
                y = height - 2 * cm
            c.drawString(x, y, ch)
            y -= line_height
        y -= 4

    for raw in md.splitlines():
        line = raw.strip()
        if not line:
            continue

        if line.startswith("#"):
            # heading
            heading = line.lstrip("#").strip()
            draw_line(heading, is_heading=True)
            continue

        if line.startswith(("-", "*")):
            draw_line("• " + line[1:].strip())
            continue

        if re.match(r"^\d+\.\s+", line):
            draw_line(line)
            continue

        draw_line(line)

    c.save()
    return buf.getvalue()


def build_zip(docx_bytes: bytes, pdf_bytes: bytes) -> bytes:
    zbuf = io.BytesIO()
    with zipfile.ZipFile(zbuf, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("notes.docx", docx_bytes)
        z.writestr("notes.pdf", pdf_bytes)
    return zbuf.getvalue()


@app.post("/make-notes")
async def make_notes(request: Request, files: List[UploadFile] = File(...)):
    _require_secret(request)

    if not files or len(files) == 0:
        raise HTTPException(status_code=400, detail="No files uploaded")
    if len(files) > MAX_FILES:
        raise HTTPException(status_code=400, detail=f"Too many files (max {MAX_FILES})")

    # Read bytes + size checks
    files_bytes: List[Tuple[str, bytes]] = []
    for f in files:
        fname = _safe_filename(f.filename or "file")
        data = _read_uploadfile(f)
        files_bytes.append((fname, data))
    _bytes_limit_check(files_bytes)

    # Extract text
    extracted_blocks = []
    for fname, data in files_bytes:
        text = extract_text(fname, data)
        if not text:
            text = "No extractable text found."
        extracted_blocks.append(f"=== File: {fname} ===\n{text}")

    source_text = "\n\n".join(extracted_blocks)

    # AI notes
    notes_md = call_ai_make_notes(source_text)

    # Build outputs
    docx_bytes = markdown_to_docx(notes_md)
    pdf_bytes = markdown_to_pdf(notes_md)

    zip_bytes = build_zip(docx_bytes, pdf_bytes)

    return StreamingResponse(
        io.BytesIO(zip_bytes),
        media_type="application/zip",
        headers={
            "Content-Disposition": 'attachment; filename="notes.zip"',
            "Cache-Control": "no-store",
        },
    )


@app.exception_handler(HTTPException)
async def http_exception_handler(request: Request, exc: HTTPException):
    return JSONResponse(status_code=exc.status_code, content={"error": exc.detail})

